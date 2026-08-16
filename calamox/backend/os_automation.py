"""OS Automation & Computer Control Engine — PyAutoGUI desktop control.

Supports:
- Native desktop software control: Notepad, WhatsApp Web/Desktop (auto-message sending), File Explorer
- File Organizer: Auto-organize Downloads by file extension
- Instant Document Generation: Auto-generate PowerPoint (.pptx) and Excel (.xlsx) spreadsheets
"""

import os
import subprocess
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    import pyautogui
    import pyperclip
    HAS_PYAUTOGUI = True
except ImportError:
    HAS_PYAUTOGUI = False


# ──────────────────────────────────────────────────────────────────────
# Helper: safe pyautogui wrapper with fail-safe awareness
# ──────────────────────────────────────────────────────────────────────
def _typewrite(text: str, pause: float = 0.01) -> None:
    """Type text character by character with small delays."""
    if HAS_PYAUTOGUI:
        pyautogui.write(text, interval=pause)
    else:
        # Fallback: type via clipboard + paste
        pyperclip.copy(text)
        pyautogui.hotkey("ctrl", "v")


def _press(key: str, presses: int = 1) -> None:
    """Press a key the specified number of times."""
    if HAS_PYAUTOGUI:
        for _ in range(presses):
            pyautogui.press(key)


# ──────────────────────────────────────────────────────────────────────
# Notepad control
# ──────────────────────────────────────────────────────────────────────
def open_notepad() -> Dict[str, Any]:
    """Launch Notepad.

    Returns:
        dict with "success": bool, "path": str|None, "error": str|None
    """
    try:
        subprocess.Popen(["notepad.exe"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        return {"success": True, "path": None, "error": None}
    except Exception as e:
        return {"success": False, "path": None, "error": str(e)}


def write_notepad(text: str, wait: float = 1.0) -> Dict[str, Any]:
    """Type text into the currently active Notepad window.

    Args:
        text: The text to type
        wait: Seconds to wait after typing

    Returns:
        dict with "success": bool, "error": str|None
    """
    try:
        if HAS_PYAUTOGUI:
            pyautogui.write(text, interval=0.01)
        else:
            pyperclip.copy(text)
            pyautogui.hotkey("ctrl", "v")
        time.sleep(wait)
        return {"success": True, "error": None}
    except Exception as e:
        return {"success": False, "error": str(e)}


def send_notepad_line(line: str) -> Dict[str, Any]:
    """Send a single line to Notepad (types text + Enter)."""
    result = write_notepad(line + "\n")
    if result["success"] and HAS_PYAUTOGUI:
        _press("enter")
    return result


# ──────────────────────────────────────────────────────────────────────
# WhatsApp Web/Desktop control
# ──────────────────────────────────────────────────────────────────────
def open_whatsapp_web() -> Dict[str, Any]:
    """Launch WhatsApp Web in the default browser.

    Returns:
        dict with "success": bool, "error": str|None
    """
    try:
        subprocess.Popen(["start", "https://web.whatsapp.com"], shell=True)
        time.sleep(3)  # Wait for page to load
        return {"success": True, "error": None}
    except Exception as e:
        return {"success": False, "error": str(e)}


def type_whatsapp_message(phone_number: str, message: str, wait_after: float = 2.0) -> Dict[str, Any]:
    """Type a message in WhatsApp Web for the given phone number.

    Assumes WhatsApp Web is already open and focused.
    Types the phone number to start a chat, then the message.

    Returns:
        dict with "success": bool, "error": str|None
    """
    try:
        # Focus address bar (Cmd+L on Mac, Ctrl+L on Windows)
        if HAS_PYAUTOGUI:
            pyautogui.hotkey("ctrl", "l")  # Focus search/address bar in most browsers
        time.sleep(0.3)

        # Type phone number (with country code format)
        formatted = phone_number.replace(" ", "").replace("-", "").replace("+", "")
        _typewrite(formatted, pause=0.02)
        time.sleep(1.0)

        # Press Enter to open chat
        if HAS_PYAUTOGUI:
            pyautogui.press("enter")
        time.sleep(wait_after)

        # Type the message
        _typewrite(message, pause=0.02)
        time.sleep(0.5)

        # Press Enter to send
        if HAS_PYAUTOGUI:
            pyautogui.press("enter")
        time.sleep(1.0)

        return {"success": True, "error": None}
    except Exception as e:
        return {"success": False, "error": str(e)}


# ──────────────────────────────────────────────────────────────────────
# File Explorer / Downloads organizer
# ──────────────────────────────────────────────────────────────────────
def organize_downloads(extensions: Optional[List[str]] = None) -> Dict[str, Any]:
    """Auto-organize Downloads folder by file extension.

    Moves files in ~/Downloads (or specified extensions) into subfolders
    named after their extension (e.g., .pdf → PDF/, .jpg → Images/).

    Args:
        extensions: List of extensions to organize (e.g., [".pdf", ".jpg"]).
                    If None, organizes ALL files.

    Returns:
        dict with "success": bool, "organized": int, "skipped": int, "error": str|None
    """
    try:
        downloads = Path.home() / "Downloads"
        if not downloads.exists():
            return {"success": False, "organized": 0, "skipped": 0, "error": f"Downloads folder not found at {downloads}"}

        if extensions is None:
            files = [f for f in downloads.iterdir() if f.is_file()]
        else:
            files = [f for f in downloads.iterdir() if f.is_file() and f.suffix.lower() in [e.lower() for e in extensions]]

        organized = 0
        skipped = 0

        for file_path in files:
            # Determine target folder name from extension
            suffix = file_path.suffix.lower()
            if not suffix:
                suffix = "unknown"

            # Map extension to folder name
            folder_map = {
                ".pdf": "PDFs",
                ".jpg": "Images",
                ".jpeg": "Images",
                ".png": "Images",
                ".gif": "Images",
                ".mov": "Videos",
                ".mp4": "Videos",
                ".mkv": "Videos",
                ".zip": "Archives",
                ".rar": "Archives",
                ".tar": "Archives",
                ".gz": "Archives",
                ".exe": "Executables",
                ".msi": "Installers",
                ".doc": "Documents",
                ".docx": "Documents",
                ".xls": "Spreadsheets",
                ".xlsx": "Spreadsheets",
                ".ppt": "Presentations",
                ".pptx": "Presentations",
                ".txt": "Textfiles",
                ".csv": "Data",
                ".py": "Scripts",
                ".js": "Scripts",
                ".html": "Web",
                ".css": "Web",
            }

            folder_name = folder_map.get(suffix, "Others")
            target_dir = downloads / folder_name
            target_dir.mkdir(parents=True, exist_ok=True)

            target_path = target_dir / file_path.name

            # Handle name conflicts
            if target_path.exists():
                stem = file_path.stem
                suffix = file_path.suffix
                counter = 1
                while target_path.exists():
                    target_path = target_dir / f"{stem}_{counter}{suffix}"
                    counter += 1

            try:
                file_path.rename(target_path)
                organized += 1
            except Exception:
                skipped += 1

        return {"success": True, "organized": organized, "skipped": skipped, "error": None}

    except Exception as e:
        return {"success": False, "organized": 0, "skipped": 0, "error": str(e)}


# ──────────────────────────────────────────────────────────────────────
# Document generation (PPTX and XLSX)
# ──────────────────────────────────────────────────────────────────────
def generate_pptx(slides: List[Dict[str, str]]) -> Tuple[bool, Optional[bytes], Optional[str]]:
    """Generate a PowerPoint (.pptx) presentation from slide data.

    Args:
        slides: List of dicts with keys "title" and "content" (or "bullet_points" as list).

    Returns:
        Tuple of (success, bytes_data, error_message)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        # Use widescreen layout
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, slide_data in enumerate(slides):
            # Add a title + content layout slide
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            title = slide.shapes.title
            title.text = slide_data.get("title", f"Slide {i + 1}")

            # Set content area
            content = slide.shapes.placeholders[1].text_frame
            content.text = ""

            data = slide_data.get("content", "") or slide_data.get("bullet_points", [])

            if isinstance(data, list):
                for item in data:
                    p = content.add_paragraph()
                    p.text = str(item)
                    p.level = 0
            elif isinstance(data, str):
                p = content.add_paragraph()
                p.text = data

            # Style the content
            for paragraph in content.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = "Calibri"

        # Save to bytes
        from io import BytesIO
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)

        return True, buf.read(), None
    except ImportError:
        return False, None, "python-pptx not installed. Install with: pip install python-pptx"
    except Exception as e:
        return False, None, str(e)


def generate_xlsx(rows: List[List[Any]], headers: Optional[List[str]] = None) -> Tuple[bool, Optional[bytes], Optional[str]]:
    """Generate an Excel (.xlsx) spreadsheet.

    Args:
        rows: List of rows, each row is a list of cell values.
        headers: Optional list of header strings (first row). If None, uses column indices.

    Returns:
        Tuple of (success, bytes_data, error_message)
    """
    try:
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"

        # Add headers if provided
        if headers:
            for col_idx, header in enumerate(headers, start=1):
                ws.cell(row=1, column=col_idx, value=header)
            data_row = 2
        elif rows and rows[0]:
            # Use first row as headers
            for col_idx, header in enumerate(rows[0], start=1):
                ws.cell(row=1, column=col_idx, value=str(header))
            data_row = 2
        else:
            data_row = 1

        # Add data rows (skip first if it was used as headers)
        data_rows = rows[1:] if (headers or (rows and rows[0])) else rows

        for row_idx, row in enumerate(data_rows, start=data_row):
            for col_idx, value in enumerate(row, start=1):
                ws.cell(row=row_idx, column=col_idx, value=value)

        # Save to bytes
        from io import BytesIO
        buf = BytesIO()
        wb.save(buf)
        buf.seek(0)

        return True, buf.read(), None
    except ImportError:
        return False, None, "openpyxl not installed. Install with: pip install openpyxl"
    except Exception as e:
        return False, None, str(e)