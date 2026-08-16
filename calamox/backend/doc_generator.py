"""Document Generator — Automated PPTX (.pptx) and XLSX (.xlsx) generation.

Generates documents via python background tasks for instant document creation.
"""

from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


# ──────────────────────────────────────────────────────────────────────
# PPTX Generation
# ──────────────────────────────────────────────────────────────────────
def generate_pptx(slides: List[Dict[str, Any]]) -> Tuple[bool, Optional[bytes], Optional[str]]:
    """Generate a PowerPoint (.pptx) presentation from slide data.

    Args:
        slides: List of dicts, each with:
            - "title" (str): Slide title
            - "content" (str | list[str]): Slide content — either a paragraph string
              or a list of bullet points.

    Returns:
        Tuple of (success, bytes_data, error_message)
        - success: True if generation succeeded
        - bytes_data: Raw pptx file bytes (for download/attachment)
        - error_message: None on success, descriptive string on failure
    """
    if not HAS_PPTX:
        return False, None, "python-pptx not installed. Install with: pip install python-pptx"

    try:
        prs = Presentation()
        # Widescreen format
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, slide_data in enumerate(slides):
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            title = slide.shapes.title
            title.text = slide_data.get("title", f"Slide {i + 1}")

            # Set content
            content_frame = slide.shapes.placeholders[1].text_frame
            content_frame.text = ""

            data = slide_data.get("content", "")
            if data is None:
                data = []

            if isinstance(data, list):
                for item in data:
                    p = content_frame.add_paragraph()
                    p.text = str(item)
                    p.level = 0
            else:
                p = content_frame.add_paragraph()
                p.text = str(data)

            # Style content text
            for paragraph in content_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = "Calibri"

        # Save to bytes buffer
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)

        return True, buf.read(), None
    except Exception as e:
        return False, None, str(e)


# ──────────────────────────────────────────────────────────────────────
# XLSX Generation
# ──────────────────────────────────────────────────────────────────────
def generate_xlsx(
    rows: List[List[Any]],
    headers: Optional[List[str]] = None,
) -> Tuple[bool, Optional[bytes], Optional[str]]:
    """Generate an Excel (.xlsx) spreadsheet.

    Args:
        rows: List of rows, each row is a list of cell values.
        headers: Optional list of header strings (first row). If None and rows
                 exist, the first row is used as headers.

    Returns:
        Tuple of (success, bytes_data, error_message)
        - success: True if generation succeeded
        - bytes_data: Raw xlsx file bytes (for download/attachment)
        - error_message: None on success, descriptive string on failure
    """
    if not HAS_OPENPYXL:
        return False, None, "openpyxl not installed. Install with: pip install openpyxl"

    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"

        # Determine header row
        if headers is not None:
            for col_idx, header in enumerate(headers, start=1):
                ws.cell(row=1, column=col_idx, value=header)
            data_row = 2
        elif rows and len(rows) > 0:
            # Use first row as headers
            for col_idx, header in enumerate(rows[0], start=1):
                ws.cell(row=1, column=col_idx, value=str(header))
            data_row = 2
        else:
            data_row = 1

        # Determine data rows
        if headers is not None:
            data_rows = rows
        elif rows and len(rows) > 0:
            data_rows = rows[1:]  # skip first row used as headers
        else:
            data_rows = rows

        for row_idx, row in enumerate(data_rows, start=data_row):
            for col_idx, value in enumerate(row, start=1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)

        # Save to bytes buffer
        buf = BytesIO()
        wb.save(buf)
        buf.seek(0)

        return True, buf.read(), None
    except Exception as e:
        return False, None, str(e)